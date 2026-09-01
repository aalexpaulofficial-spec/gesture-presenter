"""Structural analysis of an uploaded presentation using python-pptx."""

from __future__ import annotations

from typing import IO, Any

from pptx import Presentation
from pptx.util import Emu

EMU_PER_POINT = 12700


def _slide_title(slide: Any) -> str | None:
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            text = slide.shapes.title.text_frame.text.strip()
            return text or None
    except (AttributeError, ValueError):
        pass
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                return text.splitlines()[0][:120]
    return None


def _notes(slide: Any) -> str | None:
    if not slide.has_notes_slide:
        return None
    text = slide.notes_slide.notes_text_frame.text.strip()
    return text or None


def analyse_presentation(stream: IO[bytes]) -> dict[str, Any]:
    """Return slide count, native dimensions and per-slide metadata."""
    presentation = Presentation(stream)
    width = Emu(presentation.slide_width or 0) / EMU_PER_POINT
    height = Emu(presentation.slide_height or 0) / EMU_PER_POINT

    slides = []
    for index, slide in enumerate(presentation.slides):
        text_elements = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    try:
                        left = float(Emu(shape.left) / EMU_PER_POINT / float(width)) if shape.left else 0.0
                        top = float(Emu(shape.top) / EMU_PER_POINT / float(height)) if shape.top else 0.0
                        w = float(Emu(shape.width) / EMU_PER_POINT / float(width)) if shape.width else 0.0
                        h = float(Emu(shape.height) / EMU_PER_POINT / float(height)) if shape.height else 0.0
                        text_elements.append({
                            "text": text,
                            "left": left,
                            "top": top,
                            "width": w,
                            "height": h,
                        })
                    except Exception:
                        pass
        
        slides.append({
            "index": index,
            "title": _slide_title(slide),
            "shape_count": len(slide.shapes),
            "notes": _notes(slide),
            "text_elements": text_elements,
        })

    return {
        "slide_count": len(slides),
        "width_points": round(float(width), 2),
        "height_points": round(float(height), 2),
        "aspect_ratio": round(float(width) / float(height), 4) if height else 1.7778,
        "slides": slides,
    }
